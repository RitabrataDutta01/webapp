import sys, os, docx, PyPDF2, threading
from fpdf import FPDF
from werkzeug.utils import secure_filename
from pptx import Presentation
from flask import Flask, render_template, request, send_file, redirect
import openpyxl

app = Flask(__name__)

def get_base_folder():
    if getattr(sys, 'frozen', False):
        return os.path.dirname(sys.executable)  
    else:
        return os.getcwd()
    
def get_output_folder():
    base_folder = get_base_folder()
    output_folder = os.path.join(base_folder, 'output')
    os.makedirs(output_folder, exist_ok=True)
    return output_folder

def get_upload_folder():
    base_folder = get_base_folder()
    upload_folder = os.path.join(base_folder, 'uploads')
    os.makedirs(upload_folder, exist_ok=True)
    return upload_folder

AllowedExtensionWord = ['docx']
AllowedExtensionPpt = ['pptx']
AllowedExtensionImg = ['jpg', 'jpeg', 'png']
AllowedExtensionPdf = ['pdf']
AllowedExtensionExcel = ['xlsx', 'xls']

app.config['UPLOAD_FOLDER'] = get_upload_folder()
app.config['OUTPUT_FOLDER'] = get_output_folder()

def allowedfilename(filename, allowed_extensions):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in allowed_extensions

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/word2pdf', methods=['GET', 'POST'])
def word2pdf():
    if request.method == 'POST':
        if 'file' not in request.files:
            return redirect(request.url)

        file = request.files['file']
        if file.filename == '':
            return redirect(request.url)

        if file and allowedfilename(file.filename, AllowedExtensionWord):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            word = docx.Document(filepath)
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font('Arial', size=15)

            for para in word.paragraphs:
                for line in para.text.splitlines():
                    pdf.cell(200, 10, txt=line, ln=True, align='L')

            output_pdf_filename = filename.rsplit('.', 1)[0] + '.pdf'
            output_pdf_path = os.path.join(app.config['OUTPUT_FOLDER'], output_pdf_filename)
            pdf.output(output_pdf_path)
            if os.path.exists(filepath):
                os.remove(filepath)
                
            if os.path.exists(output_pdf_path):
                return send_file(output_pdf_path, as_attachment=True, download_name=output_pdf_filename)
    return render_template('word2pdf.html')

@app.route('/ppt2pdf', methods=['GET', 'POST'])
def ppt2pdf():
    if request.method == 'POST':
        if 'file' not in request.files:
            return redirect(request.url)

        file = request.files['file']
        if file.filename == '':
            return redirect(request.url)

        if file and allowedfilename(file.filename, AllowedExtensionPpt):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            ppt = Presentation(filepath)
            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=15)

            for slide in ppt.slides:
                pdf.add_page()
                pdf.set_font("Arial", size=12)

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        cleaned_text = shape.text.replace('’', "'").replace('“', '"').replace('”', '"')
                        lines = cleaned_text.split('\n')

                        for line in lines:
                            pdf.multi_cell(0, 10, line)

                pdf.ln(10)

            output_pdf_filename = filename.rsplit('.', 1)[0] + '.pdf'
            output_pdf_path = os.path.join(app.config['OUTPUT_FOLDER'], output_pdf_filename)
            pdf.output(output_pdf_path)
            if os.path.exists(filepath):
                os.remove(filepath)
                
            if os.path.exists(output_pdf_path):
                return send_file(output_pdf_path, as_attachment=True, download_name=output_pdf_filename)

    return render_template('PPT2PDF.html')

@app.route('/img2pdf', methods=['GET', 'POST'])
def img2pdf():
    if request.method == 'POST':
        if 'files' not in request.files:
            return redirect(request.url)

        files = request.files.getlist('files')
        if len(files) == 0:
            return redirect(request.url)

        file_paths = []
        for file in files:
            if file.filename == '':
                continue

            if file and allowedfilename(file.filename, AllowedExtensionImg):
                filename = secure_filename(file.filename)
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(filepath)
                file_paths.append(filepath)

            else:
                return render_template('img2pdf.html', message='Invalid file format. Please upload image files.')

        output_pdf_path = os.path.join(get_output_folder(), 'merged_images.pdf')
        convert_images_to_pdf(file_paths, output_pdf_path)

        if os.path.exists(filepath):
                os.remove(filepath)
        
        if os.path.exists(output_pdf_path):
            return send_file(output_pdf_path, as_attachment=True, download_name=output_pdf_filename)

    return render_template('img2pdf.html')

def convert_images_to_pdf(image_paths, output_pdf_path):
    pdf = FPDF()
    for image_path in image_paths:
        pdf.add_page()
        pdf.image(image_path, x=10, y=10, w=190)
    pdf.output(output_pdf_path)

@app.route('/pdfmerge', methods=['GET', 'POST'])
def pdfmerge():
    if request.method == 'POST':
        if 'files' not in request.files:
            return redirect(request.url)

        files = request.files.getlist('files')
        if len(files) == 0:
            return redirect(request.url)

        file_paths = []
        for file in files:
            if file.filename == '':
                continue

            if file and allowedfilename(file.filename, AllowedExtensionPdf):
                filename = secure_filename(file.filename)
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(filepath)
                file_paths.append(filepath)

            else:
                return render_template('pdfmerge.html', message='Invalid file format. Please upload PDF files.')

        merger = PyPDF2.PdfMerger()
        for file_path in file_paths:
            merger.append(file_path)

        output_pdf_path = os.path.join(app.config['OUTPUT_FOLDER'], 'merged.pdf')
        merger.write(output_pdf_path)
        merger.close()

        if os.path.exists(filepath):
                os.remove(filepath)

        if os.path.exists(output_pdf_path):
            return send_file(output_pdf_path, as_attachment=True, download_name=output_pdf_filename)

    return render_template('pdfmerge.html')

@app.route('/excel2pdf', methods=['GET', 'POST'])
def excel2pdf():
    if request.method == 'POST':
        if 'file' not in request.files:
            return redirect(request.url)

        file = request.files['file']
        if file.filename == '':
            return redirect(request.url)

        if file and allowedfilename(file.filename, AllowedExtensionExcel):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            wb = openpyxl.load_workbook(filepath)
            sheet = wb.active

            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)

            col_widths = []
            for col_index in range(sheet.max_column):
                max_length = 0
                for row in sheet.iter_rows(min_col=col_index + 1, max_col=col_index + 1):
                    cell_value = str(row[0].value) if row[0].value else ""
                    max_length = max(max_length, len(cell_value))

                col_widths.append(max_length * 2)

            total_width = sum(col_widths)
            page_width = 210 
            margin = 10
            available_width = page_width - margin * 2

            if total_width > available_width:
                scaling_factor = available_width / total_width
                col_widths = [width * scaling_factor for width in col_widths]

            x_pos = margin
            y_pos = 10

            for row in sheet.iter_rows():
                for col_index, cell in enumerate(row):
                    if cell.value:
                        pdf.set_xy(x_pos, y_pos)
                        pdf.cell(col_widths[col_index], 10, txt=str(cell.value), border=1, align='L')
                    x_pos += col_widths[col_index]
                x_pos = margin
                y_pos += 10

            output_pdf_filename = filename.rsplit('.', 1)[0] + '.pdf'
            output_pdf_path = os.path.join(app.config['OUTPUT_FOLDER'], output_pdf_filename)
            pdf.output(output_pdf_path)

            if os.path.exists(filepath):
                os.remove(filepath)
            
            if os.path.exists(output_pdf_path):
                return send_file(output_pdf_path, as_attachment=True, download_name=output_pdf_filename)

    return render_template('excel2pdf.html')

@app.route('/aboutme', methods=['GET'])
def aboutme():
    return render_template('aboutme.html')

def run_flask_app():
    app.run(host='0.0.0.0', port=int(os.environ.get('PORT', 5000)))
